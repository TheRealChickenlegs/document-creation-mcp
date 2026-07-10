from __future__ import annotations

import io
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

from .models import PresentationPlan, SlideSpec
from .theme_manager import Theme

# 16:9 slide dimensions.
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.7)


def _rgb(hex_color: str) -> RGBColor:
    cleaned = hex_color.lstrip("#")
    if len(cleaned) == 3:
        cleaned = "".join(c * 2 for c in cleaned)
    return RGBColor.from_string(cleaned)


def _image_dimensions(path: str) -> tuple[int, int]:
    with Image.open(path) as img:
        return img.size


def _fit_image(path: str, max_w: Emu, max_h: Emu) -> tuple[int, int]:
    """Scale image dimensions to fit within a bounding box, preserving aspect ratio."""
    w, h = _image_dimensions(path)
    ratio = min(max_w / w, max_h / h)
    return int(w * ratio), int(h * ratio)


def _set_bg(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_textbox(slide, left, top, width, height, *, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    return box, tf


def _style_run(run, size, color: RGBColor, font_name: str, *, bold=False):
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.name = font_name
    run.font.bold = bold


def _fill_bullets(tf, bullets, theme: Theme) -> None:
    body_font = theme.fonts.get("body", "Calibri")
    text_color = _rgb(theme.colors.get("text", "#222222"))
    for i, item in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(8)
        p.level = 0
        run = p.add_run()
        run.text = f"•  {item}"
        _style_run(run, 18, text_color, body_font)


def _place_image(slide, path: str, left: Emu, top: Emu, width: Emu, height: Emu) -> None:
    iw, ih = _fit_image(path, width, height)
    # Center within the requested box.
    left = left + (width - iw) // 2
    top = top + (height - ih) // 2
    slide.shapes.add_picture(path, left, top, iw, ih)


def _add_logo(slide, theme: Theme) -> None:
    logo = theme.logo
    if not logo:
        return
    try:
        size = Inches(0.6)
        _place_image(slide, logo, MARGIN, SLIDE_H - size - Inches(0.3), size, size)
    except Exception:
        # Logo is best-effort; never block deck creation.
        pass


async def _resolve_image(slide: SlideSpec, theme: Theme) -> str | None:
    from . import comfy_client

    img = slide.image
    if img is None:
        return None
    if img.source:
        return img.source
    prompt = img.prompt
    if img.use_theme_style and theme.image_style:
        prompt = f"{prompt}, {theme.image_style}"
    return await comfy_client.generate_image(
        prompt,
        negative_prompt=img.negative_prompt,
        size=img.size,
    )


def _layout_title(slide, spec: SlideSpec, theme: Theme) -> None:
    heading_font = theme.fonts.get("heading", "Calibri")
    primary = _rgb(theme.colors.get("primary", "#1F4E79"))
    text_color = _rgb(theme.colors.get("text", "#222222"))

    box, tf = _add_textbox(slide, MARGIN, Inches(2.6), SLIDE_W - 2 * MARGIN, Inches(1.6))
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    _style_run(p.add_run(), 44, primary, heading_font, bold=True)
    p.runs[0].text = spec.title

    if spec.subtitle:
        sub, stf = _add_textbox(slide, MARGIN, Inches(4.2), SLIDE_W - 2 * MARGIN, Inches(1.0))
        sp = stf.paragraphs[0]
        sp.alignment = PP_ALIGN.CENTER
        _style_run(sp.add_run(), 22, text_color, theme.fonts.get("body", "Calibri"))
        sp.runs[0].text = spec.subtitle


def _layout_section(slide, spec: SlideSpec, theme: Theme) -> None:
    accent = _rgb(theme.colors.get("accent", "#C00000"))
    heading_font = theme.fonts.get("heading", "Calibri")
    box, tf = _add_textbox(slide, MARGIN, Inches(3.0), SLIDE_W - 2 * MARGIN, Inches(1.4))
    p = tf.paragraphs[0]
    _style_run(p.add_run(), 40, accent, heading_font, bold=True)
    p.runs[0].text = spec.title


def _layout_title_content(slide, spec: SlideSpec, theme: Theme) -> None:
    heading_font = theme.fonts.get("heading", "Calibri")
    primary = _rgb(theme.colors.get("primary", "#1F4E79"))

    box, tf = _add_textbox(slide, MARGIN, MARGIN, SLIDE_W - 2 * MARGIN, Inches(1.2))
    p = tf.paragraphs[0]
    _style_run(p.add_run(), 32, primary, heading_font, bold=True)
    p.runs[0].text = spec.title

    content_top = Inches(2.0)
    if spec.bullets:
        if spec.image and spec.image.target == "content":
            # Two-column: text left, image right.
            col_w = (SLIDE_W - 3 * MARGIN) // 2
            _fill_bullets_text(slide, spec, theme, MARGIN, content_top, col_w,
                               SLIDE_H - content_top - MARGIN)
            img_path = getattr(slide, "_resolved_image", None)
            if img_path:
                _place_image(slide, img_path, MARGIN * 2 + col_w, content_top, col_w,
                             SLIDE_H - content_top - MARGIN)
        else:
            _fill_bullets_text(slide, spec, theme, MARGIN, content_top,
                               SLIDE_W - 2 * MARGIN, SLIDE_H - content_top - MARGIN)
    elif getattr(slide, "_resolved_image", None):
        img_path = slide._resolved_image
        _place_image(slide, img_path, MARGIN, content_top, SLIDE_W - 2 * MARGIN,
                     SLIDE_H - content_top - MARGIN)


def _fill_bullets_text(slide, spec, theme, left, top, width, height) -> None:
    box, tf = _add_textbox(slide, left, top, width, height, anchor=MSO_ANCHOR.TOP)
    _fill_bullets(tf, spec.bullets or [], theme)


def _layout_two_column(slide, spec: SlideSpec, theme: Theme) -> None:
    heading_font = theme.fonts.get("heading", "Calibri")
    primary = _rgb(theme.colors.get("primary", "#1F4E79"))
    box, tf = _add_textbox(slide, MARGIN, MARGIN, SLIDE_W - 2 * MARGIN, Inches(1.2))
    _style_run(tf.paragraphs[0].add_run(), 32, primary, heading_font, bold=True)
    tf.paragraphs[0].runs[0].text = spec.title

    col_w = (SLIDE_W - 3 * MARGIN) // 2
    top = Inches(2.0)
    if spec.bullets:
        _fill_bullets_text(slide, spec, theme, MARGIN, top, col_w,
                           SLIDE_H - top - MARGIN)
    img_path = getattr(slide, "_resolved_image", None)
    if img_path:
        _place_image(slide, img_path, MARGIN * 2 + col_w, top, col_w,
                     SLIDE_H - top - MARGIN)


def _layout_image_full(slide, spec: SlideSpec, theme: Theme) -> None:
    heading_font = theme.fonts.get("heading", "Calibri")
    primary = _rgb(theme.colors.get("primary", "#1F4E79"))
    img_path = getattr(slide, "_resolved_image", None)
    if img_path:
        # Full-bleed image with a title bar overlay at the top.
        _place_image(slide, img_path, Emu(0), Emu(0), SLIDE_W, SLIDE_H)
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), SLIDE_W, Inches(1.3))
        bar.fill.solid()
        bar.fill.fore_color.rgb = _rgb(theme.colors.get("background", "#000000"))
        bar.line.fill.background()
        bar.shadow.inherit = False
        tf = bar.text_frame
        tf.word_wrap = True
        _style_run(tf.paragraphs[0].add_run(), 28, primary, heading_font, bold=True)
        tf.paragraphs[0].runs[0].text = spec.title
    else:
        _layout_title_content(slide, spec, theme)


_LAYOUTS = {
    "title": _layout_title,
    "section": _layout_section,
    "title_and_content": _layout_title_content,
    "two_column": _layout_two_column,
    "image_full": _layout_image_full,
}


async def build_presentation(plan: PresentationPlan, theme, output_dir: Path) -> Path:
    """Build a .pptx from a plan + theme. Returns the saved file path."""
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    blank = prs.slide_layouts[6]
    bg_color = _rgb(theme.colors.get("background", "#FFFFFF"))

    for spec in plan.slides:
        slide = prs.slides.add_slide(blank)
        _set_bg(slide, bg_color)

        resolved = await _resolve_image(spec, theme)
        if resolved:
            slide._resolved_image = resolved  # type: ignore[attr-defined]

        layout = spec.layout
        if layout == "title_and_content" and spec.image and spec.image.target == "background":
            layout = "image_full"
        _LAYOUTS.get(layout, _layout_title_content)(slide, spec, theme)
        _add_logo(slide, theme)

    filename = plan.output_filename or plan.title.replace(" ", "_")
    if not filename.endswith(".pptx"):
        filename += ".pptx"
    out_path = output_dir / filename
    prs.save(out_path)
    return out_path
